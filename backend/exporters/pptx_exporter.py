import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

try:
    from backend.schemas import PresentationDeck
except ImportError:
    from schemas import PresentationDeck

# Color Palette
COLOR_NAVY_DARK = RGBColor(15, 23, 42)      # Slate 900
COLOR_BLUE_ACCENT = RGBColor(37, 99, 235)   # Blue 600
COLOR_TEXT_DARK = RGBColor(30, 41, 59)      # Slate 800
COLOR_TEXT_MUTED = RGBColor(148, 163, 184)  # Slate 400
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_CARD_BG = RGBColor(241, 245, 249)     # Slate 100
COLOR_CARD_BORDER = RGBColor(203, 213, 225) # Slate 300


def generate_pptx_file(deck: PresentationDeck, output_path: str = "output.pptx") -> str:
    """
    Builds a 16:9 widescreen presentation using python-pptx.
    Includes a high-impact Title Slide, structured content slides with
    bullet points and visual concept cards, and injects speaker notes.
    Returns the absolute path to the saved .pptx file.
    """
    abs_output_path = os.path.abspath(output_path)
    output_dir = os.path.dirname(abs_output_path)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)

    prs = Presentation()

    # 1. Configure 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ----------------------------------------------------
    # 2. SLIDE 1: High-Impact Title Slide
    # ----------------------------------------------------
    title_slide = prs.slides.add_slide(blank_layout)

    # Dark background fill
    bg = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_NAVY_DARK
    bg.line.fill.background()

    # Top accent bar
    top_bar = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.2)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_BLUE_ACCENT
    top_bar.line.fill.background()

    # Title & Subtitle block
    t_box = title_slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.333), Inches(3.2))
    tf = t_box.text_frame
    tf.word_wrap = True

    p_title = tf.paragraphs[0]
    p_title.text = deck.deck_title
    p_title.font.bold = True
    p_title.font.size = Pt(44)
    p_title.font.color.rgb = COLOR_WHITE
    p_title.alignment = PP_ALIGN.LEFT

    p_sub = tf.add_paragraph()
    p_sub.text = f"Target Audience: {deck.target_audience}"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_sub.space_before = Pt(18)
    p_sub.alignment = PP_ALIGN.LEFT

    # ----------------------------------------------------
    # 3. SUBSEQUENT CONTENT SLIDES
    # ----------------------------------------------------
    for slide_data in deck.slides:
        slide = prs.slides.add_slide(blank_layout)

        # Header accent indicator
        indicator = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.7), Inches(0.15), Inches(0.7)
        )
        indicator.fill.solid()
        indicator.fill.fore_color.rgb = COLOR_BLUE_ACCENT
        indicator.line.fill.background()

        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.6), Inches(11.4), Inches(0.9))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_head = tf_title.paragraphs[0]
        p_head.text = f"{slide_data.slide_num}. {slide_data.title}"
        p_head.font.size = Pt(28)
        p_head.font.bold = True
        p_head.font.color.rgb = COLOR_NAVY_DARK

        # Left Column: Bullet Points
        bullets_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(7.5), Inches(4.8))
        tf_bullets = bullets_box.text_frame
        tf_bullets.word_wrap = True

        for i, bp in enumerate(slide_data.bullet_points):
            p = tf_bullets.paragraphs[0] if i == 0 else tf_bullets.add_paragraph()
            p.text = f"•  {bp}"
            p.font.size = Pt(17)
            p.font.color.rgb = COLOR_TEXT_DARK
            p.space_before = Pt(12)

        # Right Column: Visual Diagram Concept Box
        if slide_data.visual_diagram_concept:
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(1.8), Inches(3.5), Inches(4.5)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = COLOR_CARD_BG
            card.line.color.rgb = COLOR_CARD_BORDER

            card_tf = card.text_frame
            card_tf.word_wrap = True

            p_tag = card_tf.paragraphs[0]
            p_tag.text = "💡 Visual Concept"
            p_tag.font.bold = True
            p_tag.font.size = Pt(16)
            p_tag.font.color.rgb = COLOR_BLUE_ACCENT

            p_concept = card_tf.add_paragraph()
            p_concept.text = slide_data.visual_diagram_concept
            p_concept.font.size = Pt(14)
            p_concept.font.color.rgb = COLOR_TEXT_DARK
            p_concept.space_before = Pt(10)

        # CRITICAL: Attach speaker notes directly to the slide's notes frame
        if slide_data.speaker_notes:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = slide_data.speaker_notes

    # 4. Save and return absolute path
    prs.save(abs_output_path)
    return abs_output_path
